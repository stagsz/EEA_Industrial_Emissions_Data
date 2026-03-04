from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(0, 51, 102)
ACCENT_COLOR = RGBColor(220, 53, 69)
TEXT_COLOR = RGBColor(33, 33, 33)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = TITLE_COLOR

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(0))
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(3)

    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(4)
        p.space_after = Pt(4)

# Slide 1
add_title_slide(prs, "Waste-to-Energy Market", "2024-2025 Market Intelligence & APCD Opportunities")

# Slide 2
add_content_slide(prs, "Executive Summary", [
    "GLOBAL MARKET:",
    "  * 2,869 active WtE plants | 576 MT/year capacity | 20+ year average age",
    "",
    "OPPORTUNITY:",
    "  * 7.2B - 16.2B EUR TAM for APCD retrofits (10 years)",
    "  * Peak upgrade cycle: 2024-2035 (aging asset replacement wave)",
    "",
    "GMAB FOCUS:",
    "  * Primary: Europe (regulatory-driven, premium pricing)",
    "  * Secondary: Asia (aging fleet + emerging regulation)",
])

# Slide 3 - Europe Priority
add_content_slide(prs, "EUROPE: Highest Priority Market", [
    "MARKET SIZE:",
    "  * 540 plants | 109 MT/year capacity",
    "",
    "REGULATORY URGENCY:",
    "  * EU Industrial Emissions Directive (IED) + BAT Conclusions",
    "  * Legal limit: 0.1 ng I-TEQ/Nm3",
    "  * BAT target: 0.05 ng I-TEQ/Nm3 (stricter)",
    "  * Enforcement: ACTIVE - 5M+ EUR penalties per violation",
])

# Slide 4
add_content_slide(prs, "EUROPE: Top 3 Key Markets", [
    "GERMANY (98 plants):",
    "  * 20+ years old | 98 LEADS for retrofits 2024-2030",
    "",
    "FRANCE (131 plants):",
    "  * 50-70 leads for upgrades + new projects",
    "",
    "UK (64 plants):",
    "  * Newer fleet, high profitability = capital available",
    "  * 20-30 leads for advanced APCD systems",
])

# Slide 5
add_content_slide(prs, "EUROPE: Investment & Payback", [
    "TYPICAL APCD SYSTEM (200,000 t/y plant):",
    "  * ACI + Baghouse + Temperature Control + CEM + Integration",
    "  * Total investment: 10-14M EUR",
    "",
    "ROI:",
    "  * Payback period: 1.5-3 years (compliance-driven)",
    "  * Compliance cost avoided: 2-5M EUR per violation",
    "  * Energy bonus: 3-6M EUR/year",
])

# Slide 6
add_content_slide(prs, "EUROPE: Sales Positioning", [
    "PRIMARY MESSAGE:",
    "  'Achieve BAT 0.05 ng I-TEQ/Nm3'",
    "  'Avoid 5M+ EUR penalties, maintain operational permits'",
    "",
    "TARGET DECISION MAKERS:",
    "  * Environmental/HSE Director (compliance focus)",
    "  * Plant Manager (operations)",
    "  * CFO (budget/ROI)",
])

# Slide 7
add_content_slide(prs, "ASIA: Rapid Growth, Emerging Regulation", [
    "MARKET SIZE:",
    "  * 2,224 plants (78% of global total)",
    "  * 438 MT/year capacity (76% of global)",
    "",
    "REGULATORY STATUS:",
    "  * Weak enforcement TODAY",
    "  * Tightening: China 2025+, SE Asia emerging",
    "",
    "OPPORTUNITY: Early-mover advantage (5-10 years)",
    "BEFORE regulation enforcement makes retrofits mandatory",
])

# Slide 8
add_content_slide(prs, "ASIA: Market Breakdown", [
    "JAPAN (1,000 plants):",
    "  * Oldest fleet globally (25-30 years old)",
    "  * MASSIVE aging asset replacement cycle",
    "",
    "CHINA (916 plants):",
    "  * Standards tightening 2025+",
    "  * Retrofit market LARGER than new construction",
    "",
    "SE ASIA (Vietnam, Thailand, Indonesia):",
    "  * Rapid growth (200+ projects in pipeline)",
])

# Slide 9
add_content_slide(prs, "NORTH AMERICA: Selective Opportunities", [
    "MARKET SIZE:",
    "  * 86 plants (only 3% of global total)",
    "  * USA: 71 plants | Canada: 15 plants",
    "",
    "MARKET CHARACTERISTICS:",
    "  * Mature, consolidating market",
    "  * Premium positioning required",
    "",
    "REGULATORY: EPA enforcement tightening",
])

# Slide 10
add_content_slide(prs, "Global Market Overview", [
    "TOP 10 MARKETS:",
    "1. Japan (1,000) | 2. China (916) | 3. South Korea (233)",
    "4. France (131) | 5. Germany (98) | 6. USA (71)",
    "7. UK (64) | 8. Italy (37) | 9. Sweden (36) | 10. Switzerland (29)",
    "",
    "REGIONAL TOTALS:",
    "Europe: 540 plants (19%)",
    "Asia: 2,224 plants (78%)",
])

# Slide 11
add_content_slide(prs, "5 Key Market Growth Drivers", [
    "1. LANDFILL RESTRICTIONS: Urbanization + land scarcity = more incineration",
    "",
    "2. AGING ASSET CYCLE (2024-2035): 2,869 plants at 20+ years = replacement wave",
    "",
    "3. REGULATORY TIGHTENING: EU BAT 0.05, China 2025+, EPA enforcement",
    "",
    "4. ENERGY ECONOMICS: Carbon credits + green energy improving ROI",
    "",
])

# Slide 12
add_content_slide(prs, "Technology Landscape: Flue Gas Cleaning", [
    "TRADITIONAL SYSTEMS (pre-2010):",
    "  * Basic baghouses, NO dioxin control, memory effect unaddressed",
    "",
    "MODERN APCD SYSTEMS (2010+):",
    "  * Activated Carbon Injection (ACI)",
    "  * Advanced baghouses + temperature control",
    "  * I-TEQ/TEQ continuous monitoring",
    "  * Heat recovery integration",
])

# Slide 13
add_content_slide(prs, "GMAB Competitive Advantages", [
    "DIOXIN-FIRST APPROACH:",
    "  * De novo synthesis prevention (200-400C zone control)",
    "  * Memory effect mitigation",
    "",
    "INTEGRATED SOLUTIONS:",
    "  * APCD + heat recovery (not just energy focus)",
    "",
    "REGIONAL EXPERTISE:",
    "  * EU BAT compliance knowledge (rare skillset)",
])

# Slide 14
add_content_slide(prs, "Competitive Landscape", [
    "ESTABLISHED SUPPLIERS:",
    "  Babcock & Wilcox, Metso, SPX Flow (deep relationships)",
    "",
    "CHINESE VENDORS:",
    "  Lower-cost systems, aggressive Asia pricing",
    "  Less focus on advanced dioxin control",
    "",
    "GMAB DIFFERENTIATION:",
    "  Position dioxin control as PRIMARY value driver",
])

# Slide 15
add_content_slide(prs, "Total Addressable Market (TAM)", [
    "CONSERVATIVE ESTIMATE:",
    "600-900 plants need APCD retrofits in mature markets (10 yrs)",
    "  * 12-18M EUR per APCD system",
    "  * TAM: 7.2B - 16.2B EUR (dioxin control focus)",
    "",
    "EUROPE SHARE: ~300-400 plants = 3.6-7.2B EUR",
    "ENERGY RECOVERY UPSIDE: 2-4B EUR additional",
])

# Slide 16
add_content_slide(prs, "Phased Market Entry Strategy", [
    "PHASE 1 (Q1-Q2 2025): EUROPE FOCUS",
    "  * Germany retrofit cycle | France capacity upgrades | UK advanced systems",
    "",
    "PHASE 2 (Q3-Q4 2025): ASIA EMERGENCE PREP",
    "  * Japan aging asset positioning | China regulation prep",
    "",
    "PHASE 3 (2026+): SE ASIA EARLY-MOVER",
    "  * Vietnam, Thailand, Indonesia growth markets",
])

# Slide 17
add_content_slide(prs, "Customized Value Propositions", [
    "EUROPE (PRIMARY):",
    "  'Achieve BAT 0.05 ng I-TEQ/Nm3, avoid 5M+ EUR penalties'",
    "",
    "ASIA (EMERGING):",
    "  'Prepare for tightening regulation, differentiate with advanced dioxin control'",
    "",
    "N. AMERICA (SELECTIVE):",
    "  'Stay ahead of EPA enforcement, maximize energy recovery ROI'",
])

# Slide 18
add_content_slide(prs, "Recommended Next Steps", [
    "IMMEDIATE (Q1): Lead generation focus on Europe (Germany, France, UK)",
    "",
    "SHORT-TERM (Q2): Sales messaging + competitive intelligence",
    "",
    "MEDIUM-TERM (Q3-Q4): Asia market positioning (Japan, China)",
    "",
    "ONGOING: SE Asia growth market development",
])

# Slide 19
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = ACCENT_COLOR

closing_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
text_frame = closing_text.text_frame
text_frame.word_wrap = True

p = text_frame.paragraphs[0]
p.text = "7-16B EUR Opportunity"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

text_frame.add_paragraph()
p = text_frame.paragraphs[1]
p.text = "2024-2035 Peak Upgrade Cycle"
p.font.size = Pt(40)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(20)

text_frame.add_paragraph()
p = text_frame.paragraphs[2]
p.text = "PHASE 1: EUROPE (highest urgency + premium pricing)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(30)

text_frame.add_paragraph()
p = text_frame.paragraphs[3]
p.text = "PHASE 2-3: Asia + N. America (growth + early-mover advantage)"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

prs.save('WTE_Market_Analysis_2024-2025.pptx')
print("[OK] PowerPoint created: WTE_Market_Analysis_2024-2025.pptx")
print("[OK] 19 slides - Europe-first focus")
print("[OK] Includes: Executive summary, Europe deep-dive (4 slides), Asia context, TAM sizing, phased strategy")
