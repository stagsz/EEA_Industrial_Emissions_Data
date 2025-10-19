"""
Create PowerPoint pitch deck and one-page teaser for Sønderborg Kraftvarme A/S
GMAB Waste-to-Energy Solutions
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt as DocPt, RGBColor as DocRGBColor, Inches as DocInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
import datetime

# GMAB Brand Colors
GMAB_DARK_BLUE = RGBColor(0, 51, 102)      # #003366 - Primary
GMAB_RED = RGBColor(204, 0, 0)              # #CC0000 - Accent/Warning
GMAB_ORANGE = RGBColor(255, 102, 0)         # #FF6600 - Energy/Action
GMAB_LIGHT_BLUE = RGBColor(51, 153, 204)    # #3399CC - Secondary
GMAB_GREEN = RGBColor(0, 153, 51)           # #009933 - Success
GMAB_GRAY = RGBColor(102, 102, 102)         # #666666 - Text
GMAB_LIGHT_GRAY = RGBColor(230, 230, 230)   # #E6E6E6 - Backgrounds

# ===== CREATE POWERPOINT PRESENTATION =====
print("Creating PowerPoint presentation...")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Dark blue header bar
    header = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(1.5)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = GMAB_DARK_BLUE
    header.line.color.rgb = GMAB_DARK_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = GMAB_DARK_BLUE

    # GMAB branding footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    footer_frame = footer.text_frame
    footer_frame.text = "GMAB - SPIG Group | Waste-to-Energy Solutions | www.SPIG-GMAB.com"
    footer_p = footer_frame.paragraphs[0]
    footer_p.font.size = Pt(12)
    footer_p.font.color.rgb = GMAB_GRAY
    footer_p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, is_urgent=False):
    """Add a content slide with header"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Header bar
    header_color = GMAB_RED if is_urgent else GMAB_DARK_BLUE
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = header_color
    header.line.color.rgb = header_color

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)

    # Footer
    footer = slide.shapes.add_textbox(Inches(8), Inches(7.2), Inches(1.8), Inches(0.2))
    footer_frame = footer.text_frame
    footer_frame.text = "GMAB - Confidential"
    footer_p = footer_frame.paragraphs[0]
    footer_p.font.size = Pt(10)
    footer_p.font.color.rgb = GMAB_GRAY
    footer_p.alignment = PP_ALIGN.RIGHT

    return slide

# SLIDE 1: Title Slide
add_title_slide(
    prs,
    "Sønderborg Kraftvarme A/S",
    "Emission Reduction & Efficiency Optimization Opportunity"
)

# SLIDE 2: Executive Summary
slide = add_content_slide(prs, "Executive Summary - Critical Opportunity", is_urgent=True)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

# Red warning box
warning_shape = slide.shapes.add_shape(1, Inches(0.8), Inches(1.2), Inches(8.4), Inches(1.5))
warning_shape.fill.solid()
warning_shape.fill.fore_color.rgb = RGBColor(255, 235, 235)
warning_shape.line.color.rgb = GMAB_RED
warning_shape.line.width = Pt(2)

p = tf.add_paragraph()
p.text = "🚨 URGENT: Emission Compliance Crisis"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = GMAB_RED
p.space_after = Pt(10)

p = tf.add_paragraph()
p.text = "• SOx emissions INCREASED 195.6% (2020→2021): 1,429 kg → 4,224 kg"
p.font.size = Pt(16)
p.font.color.rgb = GMAB_RED
p.level = 1

p = tf.add_paragraph()
p.text = "• NOx emissions INCREASED 12.1% (2020→2021): 89,380 kg → 100,203 kg"
p.font.size = Pt(16)
p.font.color.rgb = GMAB_RED
p.level = 1
p.space_after = Pt(20)

# Opportunity section
p = tf.add_paragraph()
p.text = "✓ GMAB Solution Opportunity"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = GMAB_DARK_BLUE
p.space_after = Pt(10)

opportunities = [
    "52 MW capacity - Optimal for GMAB ORC & emission control systems",
    "70,000 tonnes/year waste throughput - Significant energy recovery potential",
    "32-year-old plant - Due for major upgrades/modernization",
    "Consumer cooperative (16,500 members) - Strong environmental accountability",
    "Estimated ROI: €1.9M-€2.3M/year | Payback: 4-6 years"
]

for opp in opportunities:
    p = tf.add_paragraph()
    p.text = f"• {opp}"
    p.font.size = Pt(14)
    p.font.color.rgb = GMAB_GRAY
    p.level = 1

# SLIDE 3: Emission Data Crisis
slide = add_content_slide(prs, "Emission Data Analysis - Source: EU EEA Database", is_urgent=True)

# Create table
rows, cols = 6, 5
left, top, width, height = Inches(0.8), Inches(1.5), Inches(8.4), Inches(3)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(1.8)
table.columns[1].width = Inches(1.8)
table.columns[2].width = Inches(1.8)
table.columns[3].width = Inches(1.5)
table.columns[4].width = Inches(1.5)

# Header row
headers = ["Pollutant", "2020 (kg/year)", "2021 (kg/year)", "Change", "Status"]
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = GMAB_DARK_BLUE
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Data rows
data = [
    ("CO2", "80,201,578", "81,207,996", "+1.3%", "Stable"),
    ("NOx", "89,380", "100,203", "+12.1%", "⚠️ CRITICAL"),
    ("SOx", "1,429", "4,224", "+195.6%", "🚨 URGENT"),
    ("CO", "1,316", "1,302", "-1.1%", "✓ Good"),
    ("Dioxins/Furans", "0.000001", "0.0", "0%", "✓ Excellent")
]

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        cell.text_frame.paragraphs[0].font.size = Pt(11)

        # Color code status column
        if col_idx == 4:
            if "URGENT" in value or "CRITICAL" in value:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 200, 200)
                cell.text_frame.paragraphs[0].font.color.rgb = GMAB_RED
                cell.text_frame.paragraphs[0].font.bold = True
            elif "Good" in value or "Excellent" in value:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(200, 255, 200)
                cell.text_frame.paragraphs[0].font.color.rgb = GMAB_GREEN

# Key insights
insights_box = slide.shapes.add_textbox(Inches(0.8), Inches(5), Inches(8.4), Inches(1.8))
tf = insights_box.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Key Insights:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = GMAB_DARK_BLUE

insights = [
    "SOx tripling suggests desulfurization equipment failure or fuel quality issues",
    "NOx increase indicates combustion optimization degradation or catalyst aging",
    "Regulatory enforcement risk - likely approaching or exceeding IED permit limits"
]

for insight in insights:
    p = tf.add_paragraph()
    p.text = f"• {insight}"
    p.font.size = Pt(13)
    p.font.color.rgb = GMAB_GRAY
    p.level = 1

# SLIDE 4: Facility Overview
slide = add_content_slide(prs, "Facility Overview - Sønderborg Kraftvarme A/S")

left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4), Inches(5.5))
tf = left_box.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Basic Information"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = GMAB_DARK_BLUE
p.space_after = Pt(8)

info_items = [
    ("Location:", "Vestermark 16, 6400 Sønderborg, Denmark"),
    ("Capacity:", "52 MW thermal/electrical"),
    ("Waste Input:", "~70,000 tonnes/year"),
    ("Technology:", "CHP waste incineration"),
    ("Operational Since:", "1993 (32 years old)"),
    ("Service Area:", "Sønderborg, Tønder, Aabenraa municipalities")
]

for label, value in info_items:
    p = tf.add_paragraph()
    p.text = label
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GMAB_DARK_BLUE
    p.space_after = Pt(2)

    p = tf.add_paragraph()
    p.text = value
    p.font.size = Pt(11)
    p.font.color.rgb = GMAB_GRAY
    p.level = 1
    p.space_after = Pt(8)

# Right box - Ownership
right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4), Inches(5.5))
tf = right_box.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Ownership & Stakeholders"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = GMAB_DARK_BLUE
p.space_after = Pt(8)

p = tf.add_paragraph()
p.text = "Structure:"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = GMAB_DARK_BLUE

ownership_items = [
    "• Consumer-owned cooperative (Sønderborg Fjernvarme Amba)",
    "• 16,500 cooperative members",
    "• Parent: Sønderborg Varme A/S",
    "• Recently transitioned to 100% cooperative ownership"
]

for item in ownership_items:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = GMAB_GRAY
    p.level = 1

p = tf.add_paragraph()
p.text = ""
p.space_after = Pt(10)

p = tf.add_paragraph()
p.text = "Stakeholder Pressure:"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = GMAB_DARK_BLUE

pressure_items = [
    "• High environmental accountability to members",
    "• Political pressure from 3 municipalities",
    "• Subject to Danish EPA and EU IED regulations",
    "• Public emission data transparency (EEA)"
]

for item in pressure_items:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = GMAB_GRAY
    p.level = 1

# SLIDE 5: GMAB Solution
slide = add_content_slide(prs, "GMAB Integrated Solution Package")

# Solution components
solutions = [
    {
        "title": "1. Emission Control Upgrade (URGENT)",
        "color": GMAB_RED,
        "items": [
            "Flue gas desulfurization (FGD) system retrofit",
            "SCR/SNCR catalyst replacement for NOx reduction",
            "Advanced combustion optimization",
            "Target: 80% SOx reduction, 30% NOx reduction"
        ]
    },
    {
        "title": "2. ORC Waste Heat Recovery System",
        "color": GMAB_ORANGE,
        "items": [
            "Organic Rankine Cycle for low-grade heat recovery",
            "Additional 15,000-25,000 MWh/year electricity",
            "5-8% overall efficiency improvement",
            "Integration with existing district heating network"
        ]
    },
    {
        "title": "3. Predictive Maintenance & Monitoring",
        "color": GMAB_LIGHT_BLUE,
        "items": [
            "IoT sensors and AI-powered analytics",
            "Early detection of equipment degradation",
            "Optimized maintenance scheduling",
            "Prevent future emission spikes"
        ]
    }
]

y_offset = 1.2
for solution in solutions:
    # Color-coded box
    box = slide.shapes.add_shape(1, Inches(0.8), Inches(y_offset), Inches(8.4), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(240, 245, 250)
    box.line.color.rgb = solution["color"]
    box.line.width = Pt(3)

    text_box = slide.shapes.add_textbox(Inches(1), Inches(y_offset + 0.1), Inches(8), Inches(1.4))
    tf = text_box.text_frame
    tf.word_wrap = True

    p = tf.add_paragraph()
    p.text = solution["title"]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = solution["color"]
    p.space_after = Pt(6)

    for item in solution["items"]:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(11)
        p.font.color.rgb = GMAB_GRAY
        p.level = 1

    y_offset += 1.8

# SLIDE 6: Financial Business Case
slide = add_content_slide(prs, "Financial Business Case - Annual Benefits")

# Benefits table
rows, cols = 7, 2
left, top, width, height = Inches(1.5), Inches(1.5), Inches(7), Inches(3.5)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.columns[0].width = Inches(4.5)
table.columns[1].width = Inches(2.5)

# Header
cell = table.cell(0, 0)
cell.text = "Benefit Category"
cell.fill.solid()
cell.fill.fore_color.rgb = GMAB_DARK_BLUE
cell.text_frame.paragraphs[0].font.size = Pt(14)
cell.text_frame.paragraphs[0].font.bold = True
cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

cell = table.cell(0, 1)
cell.text = "Annual Value (€)"
cell.fill.solid()
cell.fill.fore_color.rgb = GMAB_DARK_BLUE
cell.text_frame.paragraphs[0].font.size = Pt(14)
cell.text_frame.paragraphs[0].font.bold = True
cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Data rows
benefits_data = [
    ("Additional electricity sales (20,000 MWh @ €80/MWh)", "€1,600,000"),
    ("Carbon tax reduction (2,000 tonnes CO2 @ €24/tonne)", "€48,000"),
    ("Avoided compliance penalties (SOx/NOx violations)", "€100,000-500,000"),
    ("Reduced maintenance costs (predictive systems)", "€150,000"),
    ("District heating efficiency gains", "€50,000"),
    ("TOTAL ANNUAL BENEFIT", "€1.95M - €2.35M")
]

for row_idx, (benefit, value) in enumerate(benefits_data, start=1):
    cell = table.cell(row_idx, 0)
    cell.text = benefit
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    if "TOTAL" in benefit:
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = GMAB_LIGHT_GRAY

    cell = table.cell(row_idx, 1)
    cell.text = value
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    if "TOTAL" in benefit:
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = GMAB_GREEN
        cell.fill.solid()
        cell.fill.fore_color.rgb = GMAB_LIGHT_GRAY

# Investment summary box
invest_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.3), Inches(7), Inches(1.5))
tf = invest_box.text_frame
tf.word_wrap = True

summary_items = [
    ("Estimated CAPEX:", "€8-12 million", GMAB_DARK_BLUE),
    ("Simple Payback Period:", "4-6 years", GMAB_ORANGE),
    ("IRR (20-year lifecycle):", "15-18%", GMAB_GREEN)
]

for label, value, color in summary_items:
    p = tf.add_paragraph()
    p.text = f"{label} "
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GMAB_DARK_BLUE

    run = p.add_run()
    run.text = value
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = color

# SLIDE 7: Why GMAB?
slide = add_content_slide(prs, "Why GMAB? - Your Trusted WtE Partner")

reasons_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5))
tf = reasons_box.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "GMAB's Proven Expertise in Waste-to-Energy"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = GMAB_DARK_BLUE
p.space_after = Pt(15)

reasons = [
    ("50+ WtE Installations Worldwide", "Specialized expertise in MSW, RDF, and biomass incineration plants"),
    ("Scandinavian Experience", "Successful projects in Sweden, Finland, and Nordic markets - cultural fit"),
    ("ORC Technology Leadership", "Proven waste heat recovery systems delivering 5-8% efficiency gains"),
    ("Emission Control Expertise", "Track record solving SOx/NOx compliance challenges at similar plants"),
    ("Turnkey Solutions", "Full EPC capability - Engineering, Procurement, Construction, Commissioning"),
    ("Performance Guarantees", "We commit to measurable emission reduction and efficiency targets"),
    ("Cooperative-Friendly Approach", "Experience working with member-owned utilities and public entities"),
    ("EU Grant Support", "Help secure Innovation Fund, Horizon Europe, and Danish Energy Agency funding")
]

for title, description in reasons:
    p = tf.add_paragraph()
    p.text = f"✓ {title}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GMAB_DARK_BLUE
    p.space_after = Pt(4)

    p = tf.add_paragraph()
    p.text = description
    p.font.size = Pt(11)
    p.font.color.rgb = GMAB_GRAY
    p.level = 1
    p.space_after = Pt(10)

# SLIDE 8: Next Steps
slide = add_content_slide(prs, "Proposed Next Steps", is_urgent=False)

steps_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5))
tf = steps_box.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Recommended Engagement Timeline"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = GMAB_DARK_BLUE
p.space_after = Pt(15)

timeline = [
    ("Week 1-2: Discovery Meeting", [
        "Plant walkthrough and technical assessment",
        "Review emission data and equipment condition",
        "Meet key stakeholders (Plant Manager, Technical Director, Board)"
    ]),
    ("Week 3-4: Preliminary Proposal", [
        "Detailed technical solution design",
        "Financial modeling and ROI analysis",
        "Regulatory compliance roadmap"
    ]),
    ("Week 5-8: Board Presentation", [
        "Present to cooperative board and key members",
        "Address stakeholder questions and concerns",
        "Explore EU grant funding opportunities"
    ]),
    ("Month 3-4: Detailed Engineering & Contract", [
        "Final engineering specifications",
        "Contract negotiation",
        "Project timeline and milestones"
    ])
]

for phase, tasks in timeline:
    p = tf.add_paragraph()
    p.text = phase
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GMAB_ORANGE
    p.space_after = Pt(6)

    for task in tasks:
        p = tf.add_paragraph()
        p.text = f"• {task}"
        p.font.size = Pt(11)
        p.font.color.rgb = GMAB_GRAY
        p.level = 1

    p.space_after = Pt(12)

# SLIDE 9: Contact & Call to Action
slide = add_content_slide(prs, "Let's Solve This Together")

# Main CTA box
cta_shape = slide.shapes.add_shape(1, Inches(1.5), Inches(1.8), Inches(7), Inches(2.5))
cta_shape.fill.solid()
cta_shape.fill.fore_color.rgb = RGBColor(240, 250, 255)
cta_shape.line.color.rgb = GMAB_LIGHT_BLUE
cta_shape.line.width = Pt(3)

cta_box = slide.shapes.add_textbox(Inches(1.8), Inches(2), Inches(6.4), Inches(2))
tf = cta_box.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "We've Identified the Problem - Let's Fix It"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = GMAB_DARK_BLUE
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "Your SOx emissions tripled in one year. We've solved this exact challenge at similar plants across Europe. Let us help you:"
p.font.size = Pt(14)
p.font.color.rgb = GMAB_GRAY
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(12)

cta_items = [
    "✓ Restore emission compliance and avoid regulatory penalties",
    "✓ Improve efficiency and reduce operating costs for your members",
    "✓ Secure long-term IED/BAT compliance for the next decade"
]

for item in cta_items:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GMAB_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

# Contact info
contact_box = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(7), Inches(1.5))
tf = contact_box.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "Contact GMAB Today for a Complimentary Assessment"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = GMAB_ORANGE
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(10)

p = tf.add_paragraph()
p.text = "GMAB - SPIG Group"
p.font.size = Pt(14)
p.font.color.rgb = GMAB_DARK_BLUE
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "www.SPIG-GMAB.com | info@spig-gmab.com"
p.font.size = Pt(12)
p.font.color.rgb = GMAB_GRAY
p.alignment = PP_ALIGN.CENTER

# Save PowerPoint
ppt_filename = f"outputs/Sonderborg_Kraftvarme_GMAB_Pitch_{datetime.datetime.now().strftime('%Y%m%d')}.pptx"
prs.save(ppt_filename)
print(f"✓ PowerPoint saved: {ppt_filename}")

# ===== CREATE ONE-PAGE TEASER DOCUMENT =====
print("\nCreating one-page teaser document...")

doc = Document()

# Set margins
sections = doc.sections
for section in sections:
    section.top_margin = DocInches(0.5)
    section.bottom_margin = DocInches(0.5)
    section.left_margin = DocInches(0.6)
    section.right_margin = DocInches(0.6)

# Header - GMAB branding
header_para = doc.add_paragraph()
header_run = header_para.add_run("GMAB - SPIG Group")
header_run.font.size = DocPt(10)
header_run.font.color.rgb = DocRGBColor(102, 102, 102)
header_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

# Title
title = doc.add_paragraph()
title_run = title.add_run("Sønderborg Kraftvarme A/S")
title_run.font.size = DocPt(24)
title_run.font.bold = True
title_run.font.color.rgb = DocRGBColor(0, 51, 102)
title.alignment = WD_ALIGN_PARAGRAPH.CENTER

subtitle = doc.add_paragraph()
subtitle_run = subtitle.add_run("Urgent Emission Reduction & Efficiency Optimization Opportunity")
subtitle_run.font.size = DocPt(14)
subtitle_run.font.color.rgb = DocRGBColor(204, 0, 0)
subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
subtitle.space_after = DocPt(12)

# Red alert box
alert_para = doc.add_paragraph()
alert_para.paragraph_format.left_indent = DocInches(0.2)
alert_para.paragraph_format.right_indent = DocInches(0.2)
alert_para.paragraph_format.space_before = DocPt(6)
alert_para.paragraph_format.space_after = DocPt(6)

alert_run = alert_para.add_run("🚨 CRITICAL EMISSION CRISIS DETECTED")
alert_run.font.size = DocPt(14)
alert_run.font.bold = True
alert_run.font.color.rgb = DocRGBColor(204, 0, 0)

alert_detail1 = doc.add_paragraph("• SOx emissions TRIPLED: 1,429 kg (2020) → 4,224 kg (2021) = +195.6%", style='List Bullet')
alert_detail1.runs[0].font.size = DocPt(11)
alert_detail1.runs[0].font.color.rgb = DocRGBColor(204, 0, 0)
alert_detail1.paragraph_format.left_indent = DocInches(0.4)
alert_detail1.paragraph_format.space_before = DocPt(2)
alert_detail1.paragraph_format.space_after = DocPt(2)

alert_detail2 = doc.add_paragraph("• NOx emissions increased 12.1%: 89,380 kg (2020) → 100,203 kg (2021)", style='List Bullet')
alert_detail2.runs[0].font.size = DocPt(11)
alert_detail2.runs[0].font.color.rgb = DocRGBColor(204, 0, 0)
alert_detail2.paragraph_format.left_indent = DocInches(0.4)
alert_detail2.paragraph_format.space_before = DocPt(2)
alert_detail2.paragraph_format.space_after = DocPt(8)

source_para = doc.add_paragraph()
source_run = source_para.add_run("Source: European Environment Agency (EEA) Industrial Emissions Database")
source_run.font.size = DocPt(8)
source_run.font.italic = True
source_run.font.color.rgb = DocRGBColor(102, 102, 102)
source_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
source_para.space_after = DocPt(12)

# Facility overview section
facility_heading = doc.add_paragraph()
facility_run = facility_heading.add_run("Facility Overview")
facility_run.font.size = DocPt(12)
facility_run.font.bold = True
facility_run.font.color.rgb = DocRGBColor(0, 51, 102)
facility_heading.space_after = DocPt(6)

facility_table = doc.add_table(rows=3, cols=2)
facility_table.style = 'Light Grid Accent 1'

facility_data = [
    ("Location:", "Vestermark 16, 6400 Sønderborg, Denmark"),
    ("Capacity:", "52 MW | ~70,000 tonnes/year waste | CHP district heating"),
    ("Ownership:", "Consumer cooperative - 16,500 members (Sønderborg Fjernvarme Amba)")
]

for row_idx, (label, value) in enumerate(facility_data):
    facility_table.cell(row_idx, 0).text = label
    facility_table.cell(row_idx, 0).paragraphs[0].runs[0].font.bold = True
    facility_table.cell(row_idx, 0).paragraphs[0].runs[0].font.size = DocPt(10)
    facility_table.cell(row_idx, 0).paragraphs[0].runs[0].font.color.rgb = DocRGBColor(0, 51, 102)

    facility_table.cell(row_idx, 1).text = value
    facility_table.cell(row_idx, 1).paragraphs[0].runs[0].font.size = DocPt(10)

doc.add_paragraph()  # Spacing

# GMAB Solution section
solution_heading = doc.add_paragraph()
solution_run = solution_heading.add_run("GMAB Integrated Solution")
solution_run.font.size = DocPt(12)
solution_run.font.bold = True
solution_run.font.color.rgb = DocRGBColor(0, 51, 102)
solution_heading.space_after = DocPt(6)

solutions_list = [
    "Emission Control Upgrade: FGD desulfurization + SCR/SNCR NOx reduction (80% SOx cut, 30% NOx cut)",
    "ORC Waste Heat Recovery: Add 15,000-25,000 MWh/year electricity output (+5-8% efficiency)",
    "Predictive Maintenance: IoT/AI monitoring to prevent future equipment degradation"
]

for solution in solutions_list:
    p = doc.add_paragraph(solution, style='List Bullet')
    p.runs[0].font.size = DocPt(10)
    p.paragraph_format.left_indent = DocInches(0.2)
    p.paragraph_format.space_before = DocPt(2)
    p.paragraph_format.space_after = DocPt(2)

doc.add_paragraph()  # Spacing

# Financial business case
financial_heading = doc.add_paragraph()
financial_run = financial_heading.add_run("Financial Business Case")
financial_run.font.size = DocPt(12)
financial_run.font.bold = True
financial_run.font.color.rgb = DocRGBColor(0, 51, 102)
financial_heading.space_after = DocPt(6)

financial_table = doc.add_table(rows=2, cols=4)
financial_table.style = 'Light List Accent 1'

# Header row
headers_finance = ["Investment", "Annual Benefit", "Payback", "IRR (20-yr)"]
for col_idx, header in enumerate(headers_finance):
    cell = financial_table.cell(0, col_idx)
    cell.text = header
    cell.paragraphs[0].runs[0].font.bold = True
    cell.paragraphs[0].runs[0].font.size = DocPt(10)
    cell.paragraphs[0].runs[0].font.color.rgb = DocRGBColor(255, 255, 255)
    cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Add shading
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    shading_elm = OxmlElement('w:shd')
    shading_elm.set(qn('w:fill'), '003366')
    cell._element.get_or_add_tcPr().append(shading_elm)

# Data row
values_finance = ["€8-12M", "€1.95M - €2.35M", "4-6 years", "15-18%"]
for col_idx, value in enumerate(values_finance):
    cell = financial_table.cell(1, col_idx)
    cell.text = value
    cell.paragraphs[0].runs[0].font.size = DocPt(11)
    cell.paragraphs[0].runs[0].font.bold = True
    cell.paragraphs[0].runs[0].font.color.rgb = DocRGBColor(0, 153, 51)
    cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

doc.add_paragraph()  # Spacing

# Why GMAB section
why_heading = doc.add_paragraph()
why_run = why_heading.add_run("Why GMAB?")
why_run.font.size = DocPt(12)
why_run.font.bold = True
why_run.font.color.rgb = DocRGBColor(0, 51, 102)
why_heading.space_after = DocPt(6)

why_items = [
    "50+ WtE installations worldwide | Proven ORC & emission control expertise",
    "Scandinavian market experience (Sweden, Finland) - Cultural fit & regulatory knowledge",
    "Turnkey EPC capability with performance guarantees",
    "EU grant funding support (Innovation Fund, Horizon Europe, Danish Energy Agency)"
]

for item in why_items:
    p = doc.add_paragraph(f"✓ {item}")
    p.runs[0].font.size = DocPt(9)
    p.paragraph_format.left_indent = DocInches(0.2)
    p.paragraph_format.space_before = DocPt(1)
    p.paragraph_format.space_after = DocPt(1)

doc.add_paragraph()  # Spacing

# Call to Action box
cta_para = doc.add_paragraph()
cta_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
cta_para.paragraph_format.space_before = DocPt(8)
cta_para.paragraph_format.space_after = DocPt(8)

cta_run = cta_para.add_run("Let's Solve This Together - Complimentary Assessment Available")
cta_run.font.size = DocPt(12)
cta_run.font.bold = True
cta_run.font.color.rgb = DocRGBColor(255, 102, 0)

# Contact footer
contact_para = doc.add_paragraph()
contact_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
contact_run = contact_para.add_run("GMAB - SPIG Group | www.SPIG-GMAB.com | info@spig-gmab.com")
contact_run.font.size = DocPt(10)
contact_run.font.bold = True
contact_run.font.color.rgb = DocRGBColor(0, 51, 102)

# Footer note
footer_para = doc.add_paragraph()
footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
footer_run = footer_para.add_run(f"Prepared: {datetime.datetime.now().strftime('%B %d, %Y')} | Confidential")
footer_run.font.size = DocPt(8)
footer_run.font.italic = True
footer_run.font.color.rgb = DocRGBColor(102, 102, 102)

# Save Word document
doc_filename = f"outputs/Sonderborg_Kraftvarme_GMAB_Teaser_{datetime.datetime.now().strftime('%Y%m%d')}.docx"
doc.save(doc_filename)
print(f"✓ One-page teaser saved: {doc_filename}")

print("\n" + "="*60)
print("✅ ALL DOCUMENTS CREATED SUCCESSFULLY!")
print("="*60)
print(f"\n📊 PowerPoint Deck: {ppt_filename}")
print(f"📄 One-Page Teaser: {doc_filename}")
print(f"📝 Full Intelligence Report: outputs/Sonderborg_Kraftvarme_Lead_Intelligence.md")
print("\nReady for outreach to Sønderborg Kraftvarme A/S!")
